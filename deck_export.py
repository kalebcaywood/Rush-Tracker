"""Build a downloadable PowerPoint of a rush day's slideshow.

Same order and content as the in-app Slideshow page: title slide, a divider
per round, one slide per PNM (latest photo + roster info), closing slide.
UT orange/white/dark styling.
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

import db

ORANGE = RGBColor(0xFF, 0x82, 0x00)
SMOKEY = RGBColor(0x58, 0x59, 0x5B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF7, 0xF7, 0xF5)

W, H = Inches(13.333), Inches(7.5)

PHOTO_BOX = (Inches(0.7), Inches(0.9), Inches(4.6), Inches(5.4))  # x, y, w, h


def _blank(prs: Presentation, bg: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def _text(slide, x, y, w, h, text, size, color, bold=False, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return box


def _labeled(slide, x, y, w, label, value, size=18):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = f"{label}   "
    r1.font.bold = True
    r1.font.color.rgb = ORANGE
    r1.font.size = Pt(size)
    r1.font.name = "Arial"
    r2 = p.add_run()
    r2.text = str(value)
    r2.font.color.rgb = DARK
    r2.font.size = Pt(size)
    r2.font.name = "Arial"


def _photo_bytes(pnm_id: str) -> bytes | None:
    photo = db.most_recent_photo(pnm_id)
    if not photo:
        return None
    try:
        return db.get_client().storage.from_(db.PHOTO_BUCKET).download(photo["storage_path"])
    except Exception:
        return None


def _add_photo(slide, data: bytes | None):
    x, y, bw, bh = PHOTO_BOX
    if data:
        try:
            from PIL import Image

            img = Image.open(io.BytesIO(data))
            iw, ih = img.size
            scale = min(bw / iw, bh / ih)
            w, h = int(iw * scale), int(ih * scale)
            slide.shapes.add_picture(
                io.BytesIO(data),
                x + int((bw - w) / 2), y + int((bh - h) / 2),
                width=w, height=h,
            )
            return
        except Exception:
            pass
    shape = slide.shapes.add_shape(1, x, y, bw, bh)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = SMOKEY
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "No photo yet"
    run.font.color.rgb = SMOKEY
    run.font.size = Pt(18)
    run.font.name = "Arial"


def build_deck(day: int) -> bytes | None:
    """Returns .pptx bytes for the day's attendance, or None if none uploaded."""
    pnms = db.attendance_for_day(day)
    if not pnms:
        return None

    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    s = _blank(prs, DARK)
    _text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
          "XI DEUTERON RUSH", 54, WHITE, bold=True)
    _text(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(0.7),
          db.DAY_LABELS.get(day, f"Day {day}"), 28, ORANGE, bold=True)

    total = len(pnms)
    seen = 0
    current_round = None
    for p in pnms:
        rnd = p.get("_round", 1)
        if rnd != current_round:
            current_round = rnd
            count = sum(1 for q in pnms if q.get("_round", 1) == rnd)
            d = _blank(prs, ORANGE)
            _text(d, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5),
                  f"ROUND {rnd}", 72, WHITE, bold=True)
            _text(d, Inches(0.9), Inches(4.3), Inches(11.5), Inches(0.6),
                  f"{count} PNMs this round", 22, DARK, bold=True)

        seen += 1
        s = _blank(prs, WHITE)
        _add_photo(s, _photo_bytes(p["id"]))
        _text(s, Inches(5.7), Inches(0.9), Inches(6.9), Inches(1.0),
              p["full_name"], 36, DARK, bold=True)
        y = 2.1
        for label, value in [
            ("Year", p.get("year")),
            ("Hometown", p.get("hometown")),
            ("High school", p.get("high_school")),
            ("RC Group", (p.get("extra") or {}).get("RC Group")),
        ]:
            if value:
                _labeled(s, Inches(5.7), Inches(y), Inches(6.9), label, value)
                y += 0.55
        notes = (p.get("notes") or "").strip()
        if notes:
            _text(s, Inches(5.7), Inches(y + 0.15), Inches(6.9), Inches(0.4),
                  "Involvement", 15, ORANGE, bold=True)
            snippet = notes[:300] + ("…" if len(notes) > 300 else "")
            _text(s, Inches(5.7), Inches(y + 0.6), Inches(6.9), Inches(1.8),
                  snippet, 13, SMOKEY)
        _text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.4),
              f"{seen} of {total}  ·  Round {rnd}  ·  Day {day}", 12, SMOKEY)

    s = _blank(prs, DARK)
    _text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.0),
          "Voting is open", 44, WHITE, bold=True)
    _text(s, Inches(0.9), Inches(3.5), Inches(11.0), Inches(1.6),
          "Open the Voting tab on your phone — your queue is in this same order. "
          "Rate every PNM before the cut meeting.", 20, RGBColor(0xCF, 0xCF, 0xCF))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
